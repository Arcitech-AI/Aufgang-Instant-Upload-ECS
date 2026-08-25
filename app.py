"""
AWS Lambda handler — single self-contained file.
Deploy this file as the Lambda package entry point.

Handler: app.handler
"""
from __future__ import annotations

# ── Standard library ──────────────────────────────────────────────────────────
import asyncio
import json
import logging
import os
import re
import sys
import time
import uuid
from datetime import datetime
from io import BytesIO, StringIO
from typing import Any, AsyncIterator, Dict, List, Optional, Tuple

# ── Third-party (lightweight only — heavy deps imported inside first use) ─────
from dotenv import load_dotenv

load_dotenv()


def _configure_lambda_writable_caches() -> None:
    """
    On AWS Lambda, $HOME and default Hugging Face / XDG cache dirs are read-only.
    Docling pulls models via huggingface_hub — point all caches at /tmp before any
    ML stack import. Optional: set HF_TOKEN in Lambda env for higher Hub rate limits.
    """
    if not os.environ.get("AWS_LAMBDA_FUNCTION_NAME"):
        return
    base = "/tmp/lambda_ml_cache"
    hf = os.path.join(base, "huggingface")
    xdg = os.path.join(base, "xdg")
    torch_home = os.path.join(base, "torch")
    for d in (hf, xdg, torch_home):
        os.makedirs(d, exist_ok=True)
    os.environ["HF_HOME"] = hf
    os.environ["HUGGINGFACE_HUB_CACHE"] = hf
    os.environ["TRANSFORMERS_CACHE"] = hf
    os.environ["XDG_CACHE_HOME"] = xdg
    os.environ["TORCH_HOME"] = torch_home


_configure_lambda_writable_caches()

# ── Logging ───────────────────────────────────────────────────────────────────
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# =============================================================================
# CONFIGURATION  (all values read from Lambda environment variables)
# =============================================================================
# OpenSearch — OPENSEARCH_HOST / OPENSEARCH_INDEX match typical .env naming;
# OPENSEARCH_URL / OPENSEARCH_INDEX_NAME remain as fallbacks for older deploys.
def _normalize_opensearch_host(raw: str) -> str:
    h = raw.strip()
    if h.startswith("https://"):
        h = h[8:]
    elif h.startswith("http://"):
        h = h[7:]
    return h.split("/")[0].split(":")[0]


OPENSEARCH_HOST = _normalize_opensearch_host(
    os.environ.get("OPENSEARCH_HOST") or os.environ.get("OPENSEARCH_URL", "")
)
OPENSEARCH_USERNAME = os.environ.get("OPENSEARCH_USERNAME", "admin").strip()
OPENSEARCH_PASSWORD = (os.environ.get("OPENSEARCH_PASSWORD") or "").strip()
OPENSEARCH_PORT = (os.environ.get("OPENSEARCH_PORT") or "443").strip()
OPENSEARCH_INDEX_NAME = (
    os.environ.get("OPENSEARCH_INDEX") or os.environ.get("OPENSEARCH_INDEX_NAME", "")
).strip()
OPENSEARCH_UNSUPPORTED_INDEX = (
    os.environ.get("OPENSEARCH_UNSUPPORTED_INDEX") or ""
).strip()

# MongoDB
MONGODB_URL = os.environ.get("MONGODB_URI", "")

# OpenAI
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")

# S3
S3_BUCKET_NAME = os.environ.get("S3_BUCKET_NAME", "bucket-aufgang-project")


def _get_param(name: str, default: Optional[str] = None) -> Optional[str]:
    """Read a config value from environment variables, with an optional default."""
    return os.environ.get(name, default)


def _require_param(name: str) -> str:
    """Like _get_param but raises RuntimeError when the value is absent or empty."""
    value = os.environ.get(name)
    if not value:
        raise RuntimeError(
            f"Required config '{name}' not found in environment."
        )
    return value

# =============================================================================
# CONFIGURATION
# =============================================================================
INSTANT_UPLOAD_COLLECTION = "instantuploadfiles"
COMPRESS_TARGET_BYTES = 20 * 1024 * 1024          # 20 MB
MAX_INSTANT_UPLOAD_CHARACTERS = 100_000
MAX_INSTANT_UPLOAD_PAGES = 10
PAGE_COUNT_EXTENSIONS = {"pdf", "docx", "pptx"}
TEXT_EXTRACT_EXTENSIONS = {"pdf", "docx", "pptx", "txt"}
OCR_LARGE_FILE_THRESHOLD_BYTES = 50 * 1024 * 1024  # 50 MB
EMBEDDING_DIM = 768

# Vectorization constants
MIN_TEXT_CONTENT_THRESHOLD = 100
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
OCR_SEMAPHORE = 6
MIN_PRINTABLE_RATIO = 0.7
MAX_TOKENS_PER_BATCH = 250_000
MAX_TOKENS_PER_INPUT = 8191
MIN_CHUNK_SIZE = 50
MAX_CONCURRENT_EMBEDDING_BATCHES = 5
EMBED_BATCH_SIZE = 50

# OCR / image-processing thresholds (aligned with vector.py)
PDF_PAGE_WEAK_TEXT_THRESHOLD = 40      # chars below which a page is treated as "weak"
IMAGE_OCR_MAX_EDGE_PX = 4000           # downscale huge rasters before Textract
IMAGE_OCR_PREPROCESS_GRAYSCALE = True  # grayscale helps Textract on prints/drawings
IMAGE_OCR_PREPROCESS_AUTOCONTRAST = True

# Docling PDF enrichment
VECTOR_PDF_DOCLING: bool = True        # enable per-page markdown enrichment via Docling
PDF_ARCH_DOCLING_PAGE_CHARS: int = 5000  # max characters kept per page from Docling output

# Textract skip codes — warn, don't crash
_TEXTRACT_SKIP_CODES = frozenset({
    "UnsupportedDocumentException",
    "InvalidParameterException",
    "BadDocumentException",
    "DocumentTooLargeException",
})

# =============================================================================
# OPTIONAL DEPS
# =============================================================================
try:
    import fitz
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    import tiktoken
    HAS_TIKTOKEN = True
except ImportError:
    HAS_TIKTOKEN = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import PyPDF2 as _pypdf2_mod

    PdfReader = _pypdf2_mod.PdfReader
    HAS_PYPDF2 = True
except ImportError:
    _pypdf2_mod = None  # noqa: F841
    PdfReader = None  # type: ignore[misc, assignment]
    HAS_PYPDF2 = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pdf2image import convert_from_bytes  # noqa: F401
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

# =============================================================================
# AWS CLIENTS  (credentials from Lambda env vars; IAM role used as fallback)
# =============================================================================
_s3_client: Optional[Any] = None
_ecs_client: Optional[Any] = None


def _get_s3_client():
    global _s3_client
    if _s3_client is None:
        import boto3

        _s3_client = boto3.client(
            "s3",
            region_name=os.environ.get("AWS_REGION", "us-east-1")
        )
    return _s3_client


def _get_ecs_client():
    global _ecs_client
    if _ecs_client is None:
        import boto3

        _ecs_client = boto3.client(
            "ecs",
            region_name=os.environ.get("AWS_REGION", "us-east-1")
        )
    return _ecs_client


def _json_default(obj: Any) -> Any:
    """JSON serializer fallback — converts datetime/date to ISO-8601 strings."""
    if isinstance(obj, datetime):
        return obj.isoformat()
    raise TypeError(f"Object of type {type(obj).__name__} is not JSON serializable")


def _dispatch_to_ecs(record_body: str, file_id: str) -> bool:
    """
    Launch an ECS Fargate task for a large file, passing the SQS record body
    via the SQS_PAYLOAD environment variable.
    Required env vars: ECS_CLUSTER, ECS_TASK_DEFINITION, ECS_CONTAINER_NAME,
                       ECS_SUBNET_IDS (comma-separated), ECS_SECURITY_GROUP_IDS (comma-separated).
    Returns True if the task was launched successfully.
    """
    cluster = os.environ.get("ECS_CLUSTER", "")
    task_def = os.environ.get("ECS_TASK_DEFINITION", "")
    container = os.environ.get("ECS_CONTAINER_NAME", "")
    subnets = [s.strip() for s in os.environ.get("ECS_SUBNET_IDS", "").split(",") if s.strip()]
    sg_ids = [s.strip() for s in os.environ.get("ECS_SECURITY_GROUP_IDS", "").split(",") if s.strip()]

    if not all([cluster, task_def, container, subnets]):
        logger.error(
            "_dispatch_to_ecs: missing ECS config for file_id=%s "
            "(ECS_CLUSTER=%r ECS_TASK_DEFINITION=%r ECS_CONTAINER_NAME=%r ECS_SUBNET_IDS=%r)",
            file_id, cluster, task_def, container, subnets,
        )
        return False

    try:
        response = _get_ecs_client().run_task(
            cluster=cluster,
            taskDefinition=task_def,
            launchType="FARGATE",
            networkConfiguration={
                "awsvpcConfiguration": {
                    "subnets": subnets,
                    "securityGroups": sg_ids,
                    "assignPublicIp": os.environ.get("ECS_ASSIGN_PUBLIC_IP", "DISABLED"),
                }
            },
            overrides={
                "containerOverrides": [{
                    "name": container,
                    "environment": [{"name": "SQS_PAYLOAD", "value": record_body}],
                }]
            },
        )

        failures = response.get("failures", [])
        if failures:
            logger.error(
                "_dispatch_to_ecs: ECS reported failures for file_id=%s: %s",
                file_id, json.dumps(failures, default=_json_default),
            )
            return False

        tasks = response.get("tasks", [])
        task_arn = tasks[0].get("taskArn", "") if tasks else ""
        logger.info(
            "_dispatch_to_ecs: launched ECS task for file_id=%s | task_arn=%s",
            file_id, task_arn,
        )
        return True
    except Exception as e:
        logger.error("_dispatch_to_ecs: failed for file_id=%s: %s", file_id, e, exc_info=True)
        return False

# =============================================================================
# PER-INVOCATION S3 LOG HANDLER
# =============================================================================

class _S3LogHandler(logging.Handler):
    """
    Buffers every log record produced during a Lambda invocation and uploads
    the full log as a single UTF-8 text file to S3 when flush_to_s3() is
    called.  One file is generated per invocation; the key is set by the
    caller so that concurrent invocations never collide.

    Log path (configurable via LOG_S3_PREFIX env var, default "logs/lambda"):
        logs/lambda/YYYY/MM/DD/<timestamp>_<request-id>.log
    """

    _FMT = logging.Formatter(
        "%(asctime)s.%(msecs)03d %(levelname)-8s [%(name)s:%(lineno)d] %(message)s",
        datefmt="%Y-%m-%dT%H:%M:%S",
    )

    def __init__(self) -> None:
        super().__init__(level=logging.DEBUG)
        self._lines: List[str] = []

    # ── logging.Handler interface ──────────────────────────────────────────

    def emit(self, record: logging.LogRecord) -> None:
        try:
            self._lines.append(self._FMT.format(record))
        except Exception:
            self.handleError(record)

    # ── Public helpers ─────────────────────────────────────────────────────

    def reset(self) -> None:
        """Discard any buffered lines (call at the start of each invocation)."""
        self._lines.clear()

    def flush_to_s3(self, bucket: str, key: str) -> bool:
        """
        Upload buffered log lines to s3://<bucket>/<key>.
        Returns True on success, False if the upload fails or there are no lines.
        Never raises — logging must not break the main request flow.
        """
        if not self._lines:
            return False
        body = "\n".join(self._lines) + "\n"
        try:
            _get_s3_client().put_object(
                Bucket=bucket,
                Key=key,
                Body=body.encode("utf-8"),
                ContentType="text/plain; charset=utf-8",
            )
            return True
        except Exception as exc:
            # Intentionally bypass the logger here to avoid infinite recursion.
            print(
                f"[_S3LogHandler] WARNING: failed to upload log to "
                f"s3://{bucket}/{key}: {exc}",
                file=sys.stderr,
            )
            return False


# Attach to the root logger so every library log is captured as well.
_s3_log_handler = _S3LogHandler()
logging.getLogger().addHandler(_s3_log_handler)


# =============================================================================
# MODULE-LEVEL MONGO (reused across warm invocations)
# =============================================================================
_mongo_client: Optional[Any] = None
_db = None


def _get_db():
    global _mongo_client, _db
    if _db is None:
        from pymongo import MongoClient

        _mongo_client = MongoClient(MONGODB_URL, serverSelectionTimeoutMS=5000)
        _db = _mongo_client.get_default_database()
    return _db


def _find_one_by_file_id(db_col, file_id: str) -> Optional[Dict[str, Any]]:
    """
    Load a file row whether ``_id`` is stored as ObjectId or as a string with
    the same 24-hex shape (common when IDs pass through JSON as plain strings).
    """
    from bson import ObjectId
    from bson.errors import InvalidId

    try:
        oid = ObjectId(file_id)
        doc = db_col.find_one({"_id": oid})
        if doc is not None:
            return doc
    except InvalidId:
        pass
    return db_col.find_one({"_id": file_id})


# Lazy-loaded Pillow (large import) — used by OCR / image paths only.
_pil_mod: Optional[Tuple[Any, Any]] = None


def _pil() -> Tuple[Any, Any]:
    global _pil_mod
    if _pil_mod is None:
        from PIL import Image as PILImage
        from PIL import ImageOps as PILImageOps

        _pil_mod = (PILImage, PILImageOps)
    return _pil_mod


# =============================================================================
# TOKENIZER HELPERS
# =============================================================================
_tokenizer = None


def _get_tokenizer():
    global _tokenizer
    if _tokenizer is None and HAS_TIKTOKEN:
        try:
            _tokenizer = tiktoken.encoding_for_model("text-embedding-3-small")
        except KeyError:
            _tokenizer = tiktoken.get_encoding("cl100k_base")
    return _tokenizer


def _count_tokens(text: str) -> int:
    tok = _get_tokenizer()
    return len(tok.encode(text)) if tok else len(text) // 4


def _chunk_validator(chunk: str) -> bool:
    return len(chunk.strip()) > MIN_CHUNK_SIZE


def _truncate_chunk(chunk: str, max_tokens: int = MAX_TOKENS_PER_INPUT) -> str:
    tok = _get_tokenizer()
    if not tok:
        max_chars = max_tokens * 4
        return chunk[:max_chars] if len(chunk) > max_chars else chunk
    tokens = tok.encode(chunk)
    if len(tokens) <= max_tokens:
        return chunk
    return tok.decode(tokens[:max_tokens])


def _create_token_batches(chunks: List[str]) -> List[List[str]]:
    batches, current, current_tokens = [], [], 0
    for chunk in chunks:
        if not _chunk_validator(chunk):
            continue
        chunk = _truncate_chunk(chunk)
        t = _count_tokens(chunk)
        if current_tokens + t > MAX_TOKENS_PER_BATCH and current:
            batches.append(current)
            current, current_tokens = [chunk], t
        else:
            current.append(chunk)
            current_tokens += t
    if current:
        batches.append(current)
    return batches


def _is_text_meaningful(text: str) -> bool:
    if not text or not text.strip():
        return False
    printable = sum(1 for c in text if c.isprintable() or c in "\n\r\t")
    ratio = printable / len(text) if text else 0
    readable = "".join(c for c in text if c.isprintable() or c in "\n\r\t").strip()
    return ratio >= MIN_PRINTABLE_RATIO and len(readable) >= MIN_TEXT_CONTENT_THRESHOLD


def _is_antiword_available() -> bool:
    import shutil

    return shutil.which("antiword") is not None


def _pdf_has_text_sync(path: str) -> Tuple[bool, int]:
    if not HAS_PYMUPDF:
        return False, 0
    doc = None
    try:
        doc = fitz.open(path)
        total = sum(len(doc[i].get_text("text").strip()) for i in range(doc.page_count))
        return total >= MIN_TEXT_CONTENT_THRESHOLD, total
    except Exception as e:
        logger.warning("_pdf_has_text_sync failed: %s", e)
        return False, 0
    finally:
        if doc:
            doc.close()


def _log_textract_failure(context: str, err: Exception) -> None:
    from botocore.exceptions import ClientError

    if isinstance(err, ClientError):
        code = (err.response or {}).get("Error", {}).get("Code", "")
        if code in _TEXTRACT_SKIP_CODES:
            msg = (err.response or {}).get("Error", {}).get("Message", str(err))
            logger.warning("Textract %s skipped (%s): %s", context, code, msg)
            return
    logger.exception("Textract %s failed: %s", context, err)


# =============================================================================
# TEXT QUALITY HELPERS  (ported from vector.py)
# =============================================================================

def normalize_text(s: str) -> str:
    """Collapse runs of spaces/tabs, strip null bytes, deduplicate blank lines."""
    s = s.replace("\x00", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{2,}", "\n", s)
    return s.strip()


def _text_symbol_noise_ratio(s: str) -> float:
    """Fraction of characters that are neither alphanumeric, space, nor common punctuation."""
    if not s:
        return 1.0
    ok = sum(
        1 for c in s
        if c.isalnum() or c.isspace() or c in ".,;:'\"°_-–()[]/\\#+*="
    )
    return 1.0 - (ok / len(s))


def is_native_text_probably_weak(
    text: str,
    char_threshold: int = PDF_PAGE_WEAK_TEXT_THRESHOLD,
) -> bool:
    """
    Return True when native-extracted text looks like a scan/drawing page that
    would benefit from OCR.  Considers: length, word/line structure, symbol noise.
    """
    s = (text or "").strip()
    if not s:
        return True
    n = len(s)
    lines = [ln.strip() for ln in s.splitlines() if ln.strip()]
    n_lines = len(lines)
    unique_lines = len({x.casefold() for x in lines})
    words = re.findall(r"\w+", s, re.UNICODE)
    n_words = len(words)
    noise = _text_symbol_noise_ratio(s)
    technical_hint = bool(re.search(r"\d", s)) and (unique_lines >= 2 or n_words >= 5)

    if n >= char_threshold * 3 and noise > 0.74:
        return True
    if n >= 400 and noise > 0.62:
        return True
    if n < char_threshold:
        if technical_hint and noise < 0.52:
            return False
        if unique_lines >= 2 and n_words >= 4 and noise < 0.58:
            return False
        return True
    if n < char_threshold * 2:
        if unique_lines >= 3 or (technical_hint and noise < 0.48):
            return False
        if noise > 0.64:
            return True
        if unique_lines <= 1 and n_words < 6:
            return True
        return False
    if noise > 0.78:
        return True
    return False


def score_extraction_candidate(text: str) -> float:
    """
    Score a text candidate so the best of native vs OCR can be selected.
    Higher = more readable / structured.
    """
    import math

    if not text or not text.strip():
        return -1e9
    s = text.strip()
    n = len(s)
    lines = [ln.strip() for ln in s.splitlines() if ln.strip()]
    n_lines = len(lines)
    n_unique_lines = len({x.casefold() for x in lines})
    printable = sum(1 for c in s if c.isprintable() or c.isspace())
    printable_ratio = printable / max(n, 1)
    alnum = sum(1 for c in s if c.isalnum())
    alnum_ratio = alnum / max(n, 1)
    tokens = re.findall(r"\S+", s)
    if not tokens:
        return -1e8
    freq: Dict[str, int] = {}
    for t in tokens:
        freq[t] = freq.get(t, 0) + 1
    repeat_penalty = sum(max(0, cnt - 2) * 1.5 for cnt in freq.values())
    if len(freq) < 4 and n > 120:
        repeat_penalty += 25.0
    symbol_like = sum(1 for c in s if not (c.isalnum() or c.isspace()))
    symbol_ratio = symbol_like / max(n, 1)
    garbage_penalty = max(0.0, (symbol_ratio - 0.35) * 120.0) if n > 25 else 0.0
    line_diversity = n_unique_lines / max(n_lines, 1)
    score = (
        min(math.log1p(n) * 18.0, 130.0)
        + n_unique_lines * 3.5
        + line_diversity * 28.0
        + printable_ratio * 42.0
        + alnum_ratio * 22.0
        - min(repeat_penalty, 95.0)
        - garbage_penalty
    )
    return score


def _preprocess_pil_image_for_ocr(pil_img: Any) -> Any:
    """Resize very large images; grayscale + autocontrast for OCR-friendly contrast."""
    Image, ImageOps = _pil()
    w, h = pil_img.size
    m = max(w, h)
    if m > IMAGE_OCR_MAX_EDGE_PX:
        scale = IMAGE_OCR_MAX_EDGE_PX / m
        nw = max(1, int(w * scale))
        nh = max(1, int(h * scale))
        pil_img = pil_img.resize((nw, nh), Image.Resampling.LANCZOS)
    if pil_img.mode not in ("L", "RGB", "RGBA"):
        pil_img = pil_img.convert("RGB")
    if IMAGE_OCR_PREPROCESS_GRAYSCALE:
        pil_img = pil_img.convert("L")
    if IMAGE_OCR_PREPROCESS_AUTOCONTRAST:
        pil_img = ImageOps.autocontrast(pil_img, cutoff=2)
    return pil_img.convert("RGB")


# =============================================================================
# PAGE / CHARACTER COUNTING  (upload_utils equivalents)
# =============================================================================

def _count_file_pages(
    file_ext: str,
    *,
    content: Optional[bytes] = None,
    path: Optional[str] = None,
) -> Optional[int]:
    """Count pages. Exactly one of path / content must be provided."""
    if (content is None) == (path is None):
        return None
    file_ext = file_ext.lower().strip(".")
    try:
        if file_ext == "pdf":
            if HAS_PYPDF2:
                if path is not None:
                    with open(path, "rb") as f:
                        return len(PdfReader(f).pages)
                return len(PdfReader(BytesIO(content)).pages)  # type: ignore[arg-type]
            if HAS_PDFPLUMBER:
                pdf_cm = (
                    pdfplumber.open(path)
                    if path is not None
                    else pdfplumber.open(BytesIO(content))  # type: ignore[arg-type]
                )
                with pdf_cm as pdf:
                    return len(pdf.pages)
            return None
        if file_ext == "docx":
            if not HAS_DOCX:
                return None
            doc = DocxDocument(path) if path is not None else DocxDocument(BytesIO(content))  # type: ignore[arg-type]
            para_count = len(doc.paragraphs)
            return max(1, para_count // 25) if para_count > 0 else 1
        if file_ext == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(path) if path is not None else Presentation(BytesIO(content))  # type: ignore[arg-type]
                return len(prs.slides)
            except ImportError:
                return None
    except Exception as e:
        logger.warning("Could not count pages for %s: %s", file_ext, e)
    return None


def _count_file_characters(
    file_ext: str,
    path: Optional[str] = None,
    file_bytes: Optional[bytes] = None,
) -> Optional[int]:
    ext = file_ext.lower().strip(".")
    if ext in ("txt", "text", "csv"):
        try:
            raw = open(path, "rb").read() if path else (file_bytes or b"")
            return len(raw.decode("utf-8", errors="ignore"))
        except OSError:
            return None
    if ext == "pdf":
        if not HAS_PYMUPDF:
            return None
        doc = None
        try:
            doc = fitz.open(path) if path else fitz.open(stream=file_bytes, filetype="pdf")
            return sum(len((doc[i].get_text("text") or "").strip()) for i in range(doc.page_count))
        except Exception:
            return None
        finally:
            if doc is not None:
                doc.close()
    if ext == "docx" and HAS_DOCX:
        if path is None and file_bytes is None:
            return None
        try:
            d = DocxDocument(path) if path is not None else DocxDocument(BytesIO(file_bytes))  # type: ignore[arg-type]
            return sum(len(p.text) for p in d.paragraphs)
        except Exception:
            return None
    return None


# =============================================================================
# COMPRESSION HELPERS
# =============================================================================

def _needs_compression(local_path: str, file_ext: str) -> bool:
    if file_ext.lower() != "pdf":
        return False
    try:
        return os.path.getsize(local_path) > COMPRESS_TARGET_BYTES
    except OSError:
        return False


def _compress_pdf(src: str, dst: str) -> bool:
    """Write recompressed PDF to dst. Returns True if dst is smaller than src."""
    if not HAS_PYMUPDF:
        logger.warning("PyMuPDF not available – cannot compress PDF.")
        return False
    doc = None
    try:
        doc = fitz.open(src)
        doc.save(dst, garbage=4, deflate=True, clean=True)
        doc.close()
        doc = None
        if not os.path.isfile(dst):
            return False
        if os.path.getsize(dst) >= os.path.getsize(src):
            _unlink_quiet(dst)
            return False
        return True
    except Exception as e:
        logger.warning("PDF compression failed: %s", e)
        _unlink_quiet(dst)
        return False
    finally:
        if doc is not None:
            doc.close()


def compress_file_if_needed(local_path: str, file_ext: str, s3_key: str, file_id: str) -> str:
    if not _needs_compression(local_path, file_ext):
        return local_path

    try:
        file_size = os.path.getsize(local_path)
    except OSError:
        return local_path

    try:
        import psutil
        available_mb = psutil.virtual_memory().available / (1024 * 1024)
        required_mb = (file_size / (1024 * 1024)) * 2
        if available_mb < required_mb:
            logger.warning(
                "compress_file_if_needed: low memory (%s MB avail, ~%s MB needed). Skipping.",
                f"{available_mb:.0f}", f"{required_mb:.0f}",
            )
            return local_path
    except ImportError:
        pass

    import tempfile

    suffix = f".{file_ext}" if file_ext else ".bin"
    fd, tmp_compressed = tempfile.mkstemp(prefix="cmp_", suffix=suffix)
    os.close(fd)
    try:
        if not _compress_pdf(local_path, tmp_compressed):
            return local_path

        original_mb = file_size / (1024 * 1024)
        os.replace(tmp_compressed, local_path)
        tmp_compressed = None
        new_size = os.path.getsize(local_path)
        logger.info(
            "Compressed %s (%s) %.1f MB → %.1f MB",
            file_id, file_ext.upper(), original_mb, new_size / (1024 * 1024),
        )
        import mimetypes

        ct = mimetypes.guess_type(file_id)[0] or "application/octet-stream"
        try:
            _get_s3_client().upload_file(
                local_path,
                S3_BUCKET_NAME,
                s3_key,
                ExtraArgs={"ContentType": ct},
            )
            logger.info("Replaced S3 object %s with compressed version.", s3_key)
        except Exception as s3_err:
            logger.warning("Failed to replace S3 with compressed: %s", s3_err)
    except Exception as comp_err:
        logger.warning("Compression failed for %s: %s", file_id, comp_err)
    finally:
        if tmp_compressed and os.path.exists(tmp_compressed):
            _unlink_quiet(tmp_compressed)
    return local_path


# =============================================================================
# S3 HELPERS
# =============================================================================

def download_s3_to_temp(s3_key: str, file_ext: str) -> str:
    import tempfile

    suffix = f".{file_ext.lstrip('.')}" if file_ext else ".bin"
    fd, path = tempfile.mkstemp(prefix="ocr_", suffix=suffix)
    os.close(fd)
    bucket = S3_BUCKET_NAME
    try:
        _get_s3_client().download_file(bucket, s3_key, path)
        return path
    except Exception:
        _unlink_quiet(path)
        raise


def _unlink_quiet(path: Optional[str]) -> None:
    if path:
        try:
            os.unlink(path)
        except OSError:
            pass


# =============================================================================
# VALIDATION
# =============================================================================

def validate_page_count(file_ext: str, path: str) -> Optional[str]:
    if file_ext not in PAGE_COUNT_EXTENSIONS:
        return None
    page_count = _count_file_pages(file_ext, path=path)
    if page_count is not None and page_count > MAX_INSTANT_UPLOAD_PAGES:
        return f"File has {page_count} pages; max is {MAX_INSTANT_UPLOAD_PAGES} for instant upload."
    return None


def validate_char_count(file_ext: str, path: str) -> Optional[str]:
    if file_ext not in TEXT_EXTRACT_EXTENSIONS:
        return None
    char_count = _count_file_characters(file_ext, path=path)
    if char_count is not None and char_count > MAX_INSTANT_UPLOAD_CHARACTERS:
        return (
            f"File has {char_count:,} characters; max is "
            f"{MAX_INSTANT_UPLOAD_CHARACTERS:,} for instant upload."
        )
    return None


# =============================================================================
# STANDALONE OPENSEARCH SERVICE
# =============================================================================

class StandaloneOpenSearchService:

    def __init__(self) -> None:
        from opensearchpy import OpenSearch, helpers

        self._bulk = helpers.bulk
        self.index_name = OPENSEARCH_INDEX_NAME
        self.unsupported_index = OPENSEARCH_UNSUPPORTED_INDEX

        _port = int(OPENSEARCH_PORT or "443")
        _use_ssl = _port == 443
        _http_auth = None
        if OPENSEARCH_PASSWORD:
            _http_auth = (OPENSEARCH_USERNAME or "admin", OPENSEARCH_PASSWORD)

        self.client = OpenSearch(
            hosts=[{"host": OPENSEARCH_HOST, "port": _port}],
            http_auth=_http_auth,
            use_ssl=_use_ssl,
            verify_certs=_use_ssl,
            timeout=60,
            max_retries=3,
            retry_on_timeout=True,
            http_compress=True,
        )

    async def is_file_indexed(self, file_id: str, file_path: str) -> bool:
        try:
            response = await asyncio.to_thread(
                self.client.count,
                index=self.index_name,
                body={"query": {"term": {"file_id": file_id}}},
            )
            return response.get("count", 0) > 0
        except Exception as e:
            logger.warning("is_file_indexed check failed for %s: %s", file_id, e)
            return False

    async def index_file(
        self,
        file_id: str,
        file_doc: Dict[str, Any],
        chunks: List[Dict[str, Any]],
        replace: bool = False,
    ):
        if not chunks:
            logger.warning("No chunks to index for file %s", file_id)
            return
        actions = []
        for chunk in chunks:
            chunk_id = chunk.get(
                "chunk_id",
                f"{file_id}_{chunk.get('page_number', 0)}_{chunk.get('chunk_index', 0)}",
            )
            doc = {
                "file_id": file_id,
                "file_name": file_doc.get("file_name", ""),
                "file_type": file_doc.get("file_type", ""),
                "file_path": file_doc.get("full_path", ""),
                "s3_key": file_doc.get("s3_key", ""),
                "file_size": file_doc.get("file_size", 0),
                "content": chunk.get("content", ""),
                "page_number": chunk.get("page_number", 0),
                "chunk_index": chunk.get("chunk_index", 0),
                "embedding": chunk.get("embedding") or [0.0] * EMBEDDING_DIM,
                "is_access_restricted": file_doc.get("is_access_restricted", False),
                "access_permissions": file_doc.get("access_permissions", {}),
                "created_at": file_doc.get("created_at", datetime.utcnow()),
                "last_modified": file_doc.get("updated_at", datetime.utcnow()),
                "is_supported": True,
                "is_uploaded": True,
                "is_favourite": False,
                "is_deleted": file_doc.get("is_deleted", False),
                "environment": (file_doc.get("environment") or "local").strip().upper(),
            }
            if file_doc.get("upload_batch_id"):
                doc["upload_batch_id"] = file_doc["upload_batch_id"]
            if replace:
                actions.append({
                    "_op_type": "index",
                    "_index": self.index_name,
                    "_id": chunk_id,
                    "_source": doc,
                })
            else:
                actions.append({
                    "_op_type": "update",
                    "_index": self.index_name,
                    "_id": chunk_id,
                    "doc": doc,
                    "doc_as_upsert": True,
                })
        for i in range(0, len(actions), 100):
            await asyncio.to_thread(self._bulk, self.client, actions[i:i + 100])
        logger.info("Indexed %d chunks for file %s", len(actions), file_id)

    async def delete_file(self, file_id: str):
        try:
            response = await asyncio.to_thread(
                self.client.search,
                index=self.index_name,
                body={"query": {"term": {"file_id": file_id}}, "size": 10000},
            )
            hits = response.get("hits", {}).get("hits", [])
            if hits:
                actions = [
                    {"_op_type": "delete", "_index": self.index_name, "_id": h["_id"]}
                    for h in hits
                ]
                await asyncio.to_thread(self._bulk, self.client, actions)
                logger.info("Deleted %d chunks for file %s", len(actions), file_id)
        except Exception as e:
            logger.error("delete_file failed for %s: %s", file_id, e, exc_info=True)
            raise

    async def index_unsupported_file(self, file_doc: Dict[str, Any]):
        file_id = str(file_doc.get("_id"))
        doc = {
            "id": file_id,
            "file_id": file_id,
            "file_name": file_doc.get("file_name", ""),
            "file_type": file_doc.get("file_type", ""),
            "file_location": file_doc.get("full_path", ""),
            "file_size": file_doc.get("file_size", 0),
            "is_access_restricted": file_doc.get("is_access_restricted", False),
            "access_permissions": file_doc.get("access_permissions", {}),
            "created_at": file_doc.get("created_at", datetime.utcnow()),
            "last_modified": file_doc.get("updated_at", datetime.utcnow()),
            "is_supported": False,
            "is_uploaded": True,
            "is_favourite": False,
            "is_deleted": file_doc.get("is_deleted", False),
            "environment": (file_doc.get("environment") or "local").strip().upper(),
        }
        await asyncio.to_thread(
            self.client.index,
            index=self.unsupported_index,
            id=file_id,
            body=doc,
        )
        logger.info("Indexed unsupported file %s", file_id)


def _docling_available() -> bool:
    try:
        import docling  # noqa: F401
        return True
    except ImportError:
        return False


# =============================================================================
# DOCX / PPTX SHAPE HELPERS  (ported from vector.py)
# =============================================================================

def _docx_collect_headers_footers(document: Any) -> List[str]:
    """Extract header and footer text from all sections of a Word document."""
    lines: List[str] = []
    try:
        for section in document.sections:
            for label, part in (
                ("Header", section.header),
                ("Footer", section.footer),
            ):
                try:
                    for p in part.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            lines.append(f"[{label}] {t}")
                except Exception:
                    pass
            if getattr(section, "different_first_page_header_footer", False):
                try:
                    for p in section.first_page_header.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            lines.append(f"[First page header] {t}")
                    for p in section.first_page_footer.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            lines.append(f"[First page footer] {t}")
                except Exception:
                    pass
    except Exception as e:
        logger.debug("DOCX header/footer extraction skipped: %s", e)
    return lines


def _pptx_collect_shape_text(shape: Any, lines: List[str]) -> None:
    """Walk groups, tables, and text frames recursively (reading order follows shape tree)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        st = getattr(shape, "shape_type", None)
        if st == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    _pptx_collect_shape_text(child, lines)
            except Exception:
                pass
            return
    except ImportError:
        pass
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            t = (shape.text or "").strip()
            if t:
                lines.append(t)
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False) and shape.table is not None:
            for row in shape.table.rows:
                cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
                if any(cells):
                    lines.append(" | ".join(cells))
    except Exception:
        pass


def _docling_pdf_pipeline_options():
    """Memory-safe Docling pipeline options (no OCR, no tables, backend text forced)."""
    from docling.datamodel.pipeline_options import ThreadedPdfPipelineOptions
    return ThreadedPdfPipelineOptions(
        do_ocr=False,
        do_table_structure=False,
        force_backend_text=True,
        images_scale=0.5,
        layout_batch_size=1,
        ocr_batch_size=1,
        table_batch_size=1,
    )


def _extract_docling_per_page_markdown(file_bytes: bytes) -> Optional[Dict[int, str]]:
    """Run Docling on a PDF and return 1-based page number → markdown string."""
    if not VECTOR_PDF_DOCLING or not file_bytes:
        return None
    if not _docling_available():
        logger.warning("VECTOR_PDF_DOCLING=true but docling is not installed; skipping")
        return None

    from docling.datamodel.base_models import InputFormat
    from docling.document_converter import DocumentConverter, PdfFormatOption
    import tempfile
    from pathlib import Path

    tmp_path: Optional[Path] = None
    try:
        fd, name = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        tmp_path = Path(name)
        tmp_path.write_bytes(file_bytes)

        opts = _docling_pdf_pipeline_options()
        converter = DocumentConverter(
            allowed_formats=[InputFormat.PDF],
            format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)},
        )
        result = converter.convert(str(tmp_path), raises_on_error=False)
        doc = result.document
        if doc is None:
            return None

        errs = getattr(result, "errors", None) or []
        if errs:
            logger.warning(
                "Docling reported %d error(s); using per-page markdown where available",
                len(errs),
            )

        n_attr = getattr(doc, "num_pages", None)
        n = int(n_attr() if callable(n_attr) else (n_attr or 0))
        if n <= 0:
            return None

        out: Dict[int, str] = {}
        for p in range(1, n + 1):
            try:
                md = doc.export_to_markdown(page_no=p)
            except TypeError:
                # Older Docling: export whole doc and split on form-feed
                if p != 1:
                    break
                full = (doc.export_to_markdown() or "").strip()
                if "\f" in full:
                    parts = [x.strip() for x in full.split("\f") if x.strip()]
                    for i, part in enumerate(parts, start=1):
                        out[i] = part[:PDF_ARCH_DOCLING_PAGE_CHARS]
                else:
                    out[1] = full[:PDF_ARCH_DOCLING_PAGE_CHARS]
                break
            except Exception as ex:
                logger.debug("Docling export page %d: %s", p, ex)
                continue
            s = (md or "").strip()
            if s:
                out[p] = s[:PDF_ARCH_DOCLING_PAGE_CHARS]

        return out if out else None

    except Exception as e:
        logger.warning("Docling conversion failed: %s", e)
        return None
    finally:
        if tmp_path is not None:
            try:
                tmp_path.unlink(missing_ok=True)
            except OSError:
                pass



# =============================================================================
# STANDALONE VECTORIZATION SERVICE
# =============================================================================

class StandaloneVectorizationService:

    def __init__(self) -> None:
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", "! ", "? ", ", ", " "],
            length_function=len,
            is_separator_regex=False,
        )
        self.embeddings = self._init_embeddings()
        self.textract_client = self._init_textract()
        self.vision_client = self._init_vision()

    def _init_embeddings(self) -> Any:
        from langchain_openai import OpenAIEmbeddings

        return OpenAIEmbeddings(
            model="text-embedding-3-small",
            dimensions=768,
            openai_api_key=OPENAI_API_KEY,
        )

    def _init_textract(self):
        try:
            import boto3

            return boto3.client(
                "textract",
                region_name=_get_param("AWS_TEXTRACT_REGION", "us-east-1"),
                aws_access_key_id=(
                    _get_param("AWS_TEXTRACT_ACCESS_KEY")
                    or _get_param("AWS_ACCESS_KEY_ID")
                    or _get_param("AWS_ACCESS_KEY")
                ),
                aws_secret_access_key=(
                    _get_param("AWS_TEXTRACT_SECRET_KEY")
                    or _get_param("AWS_SECRET_ACCESS_KEY")
                    or _get_param("AWS_SECRET_KEY")
                ),
            )
        except Exception as e:
            logger.warning("Textract init failed: %s", e)
            return None

    def _init_vision(self) -> Optional[Any]:
        from openai import AsyncOpenAI

        api_key = OPENAI_API_KEY
        return AsyncOpenAI(api_key=api_key) if api_key else None

    # ── Public: batched vectorization ─────────────────────────────────────────

    async def vectorize_file_batched(
        self,
        file_id: str,
        file_doc: Dict[str, Any],
        file_type: str,
        file_bytes: Optional[bytes] = None,
        use_ocr: bool = True,
        file_path: Optional[str] = None,
    ) -> AsyncIterator[Dict[str, Any]]:
        _t0 = time.monotonic()
        logger.info(
            "vectorize_file_batched: START | file_id=%s | type=%s | use_ocr=%s | "
            "source=%s",
            file_id, file_type, use_ocr,
            f"path:{file_path}" if file_path else "bytes",
        )
        try:
            import gc

            # ── Text extraction ────────────────────────────────────────────
            logger.info("vectorize_file_batched: extracting text | file_id=%s | type=%s", file_id, file_type)
            _t = time.monotonic()
            if file_path and os.path.isfile(file_path):
                extracted = await self._extract_text(None, file_type, use_ocr=use_ocr, file_path=file_path)
            elif file_bytes is not None:
                extracted = await self._extract_text(file_bytes, file_type, use_ocr=use_ocr)
            else:
                logger.warning("vectorize_file_batched: no content provided for file_id=%s", file_id)
                return

            if not extracted:
                logger.warning(
                    "vectorize_file_batched: text extraction returned nothing | file_id=%s | type=%s | elapsed=%.3fs",
                    file_id, file_type, time.monotonic() - _t,
                )
                return

            _total_chars = sum(len(t) for t, _ in extracted)
            _pages = len(extracted)
            logger.info(
                "vectorize_file_batched: extraction complete | file_id=%s | pages=%d | "
                "total_chars=%d | elapsed=%.3fs",
                file_id, _pages, _total_chars, time.monotonic() - _t,
            )

            # ── Chunking ───────────────────────────────────────────────────
            logger.info("vectorize_file_batched: chunking content | file_id=%s", file_id)
            _t = time.monotonic()
            chunks = await self._chunk_content(extracted, file_id)
            if not chunks:
                logger.warning(
                    "vectorize_file_batched: chunking produced 0 chunks | file_id=%s | elapsed=%.3fs",
                    file_id, time.monotonic() - _t,
                )
                return
            logger.info(
                "vectorize_file_batched: chunking complete | file_id=%s | chunks=%d | elapsed=%.3fs",
                file_id, len(chunks), time.monotonic() - _t,
            )

            # ── Embed + yield ──────────────────────────────────────────────
            total_batches = (len(chunks) + EMBED_BATCH_SIZE - 1) // EMBED_BATCH_SIZE
            for start in range(0, len(chunks), EMBED_BATCH_SIZE):
                batch_idx = start // EMBED_BATCH_SIZE + 1
                batch = chunks[start:start + EMBED_BATCH_SIZE]
                logger.info(
                    "vectorize_file_batched: embedding batch %d/%d | size=%d | file_id=%s",
                    batch_idx, total_batches, len(batch), file_id,
                )
                _t = time.monotonic()
                embeddings = await self._generate_embeddings(batch)
                if not embeddings:
                    logger.warning(
                        "vectorize_file_batched: no embeddings returned for batch %d/%d | file_id=%s",
                        batch_idx, total_batches, file_id,
                    )
                    continue
                logger.info(
                    "vectorize_file_batched: embeddings received | batch=%d/%d | count=%d | elapsed=%.3fs",
                    batch_idx, total_batches, len(embeddings), time.monotonic() - _t,
                )
                min_len = min(len(embeddings), len(batch))
                result_chunks = []
                null_emb = 0
                for i, (chunk, emb) in enumerate(zip(batch[:min_len], embeddings[:min_len])):
                    if emb is not None:
                        result_chunks.append({
                            "chunk_id": f"{file_id}_{chunk.get('page_number', 0)}_{start + i}",
                            "chunk_index": start + i,
                            "page_number": chunk.get("page_number", 0),
                            "content": chunk.get("content", ""),
                            "embedding": emb,
                        })
                    else:
                        null_emb += 1
                if null_emb:
                    logger.warning(
                        "vectorize_file_batched: %d null embeddings in batch %d/%d | file_id=%s",
                        null_emb, batch_idx, total_batches, file_id,
                    )
                if result_chunks:
                    logger.debug(
                        "vectorize_file_batched: yielding %d chunks (batch %d/%d) | file_id=%s",
                        len(result_chunks), batch_idx, total_batches, file_id,
                    )
                    yield {"chunks": result_chunks}
                gc.collect()

            logger.info(
                "vectorize_file_batched: DONE | file_id=%s | total_elapsed=%.3fs",
                file_id, time.monotonic() - _t0,
            )
        except Exception as e:
            logger.error(
                "vectorize_file_batched: UNHANDLED ERROR | file_id=%s | elapsed=%.3fs | error=%s",
                file_id, time.monotonic() - _t0, e,
                exc_info=True,
            )

    # ── Text extraction ────────────────────────────────────────────────────────

    async def _extract_text(
        self,
        file_bytes: Optional[bytes],
        file_type: str,
        use_ocr: bool = False,
        file_path: Optional[str] = None,
    ) -> Optional[List[Tuple[str, int]]]:
        use_path = bool(file_path and os.path.isfile(file_path))
        try:
            if file_type == "pdf":
                return await self._extract_pdf_text(file_bytes=file_bytes, use_ocr=use_ocr, file_path=file_path if use_path else None)
            elif file_type == "docx":
                return await self._extract_docx_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type == "doc":
                return await self._extract_doc_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type in ("xlsx", "xls", "xlsm", "excel"):
                return await self._extract_excel_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type in ("pptx", "ppt"):
                return await self._extract_ppt_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type == "csv":
                return await self._extract_csv_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type in ("txt", "text"):
                if use_path:
                    text = await asyncio.to_thread(lambda: open(file_path, "r", encoding="utf-8", errors="ignore").read())
                else:
                    text = (file_bytes or b"").decode("utf-8", errors="ignore")
                return [(normalize_text(text), 1)] if text.strip() else None
            elif file_type in ("png", "jpg", "jpeg", "jpe", "img", "image", "gif", "bmp", "webp"):
                return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type in ("tiff", "tif"):
                return await self._extract_tiff_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            elif file_type == "dat":
                raw = await asyncio.to_thread(lambda: open(file_path, "rb").read()) if use_path else (file_bytes or b"")
                text = raw.decode("utf-8", errors="ignore")
                return [(text, 1)] if text.strip() else None
            elif file_type == "bin":
                return await self._extract_bin_text(file_bytes=file_bytes, file_path=file_path if use_path else None)
            else:
                logger.warning("Unsupported file type: %s", file_type)
                return None
        except Exception as e:
            logger.error("_extract_text error: %s", e, exc_info=True)
            return None

    # ── PDF ────────────────────────────────────────────────────────────────────

    def _extract_pdf_pages_sync(self, path: str) -> List[Tuple[str, int]]:
        doc = None
        results = []
        try:
            doc = fitz.open(path)
            if doc.is_encrypted:
                raise ValueError("PDF is encrypted")
            for i in range(doc.page_count):
                text = doc[i].get_text("text").strip()
                if text:
                    results.append((text, i + 1))
            return results
        finally:
            if doc:
                doc.close()
            import gc

            gc.collect()

    async def _extract_pdf_text(
        self,
        file_bytes: Optional[bytes] = None,
        use_ocr: bool = False,
        file_path: Optional[str] = None,
    ) -> Optional[List[Tuple[str, int]]]:
        if not file_bytes and not file_path:
            return None

        cached_bytes: Optional[bytes] = None

        async def _bytes() -> bytes:
            nonlocal cached_bytes
            if cached_bytes is None:
                if file_bytes is not None:
                    cached_bytes = file_bytes
                else:
                    cached_bytes = await asyncio.to_thread(lambda: open(file_path, "rb").read())
            return cached_bytes

        if file_path:
            has_text, total_chars = await asyncio.to_thread(_pdf_has_text_sync, file_path)
        else:
            has_text, total_chars = False, 0
            if HAS_PYMUPDF:
                b = await _bytes()
                def _check(b):
                    doc = fitz.open(stream=b, filetype="pdf")
                    t = sum(len(doc[i].get_text("text").strip()) for i in range(doc.page_count))
                    doc.close()
                    return t >= MIN_TEXT_CONTENT_THRESHOLD, t
                has_text, total_chars = await asyncio.to_thread(_check, b)

        page_texts: List[Tuple[str, int]] = []
        try:
            if HAS_PYMUPDF and file_path:
                page_texts = await asyncio.to_thread(self._extract_pdf_pages_sync, file_path)
            elif HAS_PYMUPDF:
                b = await _bytes()
                def _extract_bytes(b):
                    doc = fitz.open(stream=b, filetype="pdf")
                    results = []
                    for i in range(doc.page_count):
                        t = doc[i].get_text("text").strip()
                        if t:
                            results.append((t, i + 1))
                    doc.close()
                    return results
                page_texts = await asyncio.to_thread(_extract_bytes, b)
            elif HAS_PYPDF2 and PdfReader is not None:
                src = file_path if file_path else BytesIO(await _bytes())
                reader = await asyncio.to_thread(PdfReader, src)
                for i, page in enumerate(reader.pages, 1):
                    t = (await asyncio.to_thread(page.extract_text) or "").strip()
                    if t:
                        page_texts.append((t, i))
            else:
                logger.warning(
                    "_extract_pdf_text: no PyMuPDF/PyPDF2 — cannot extract native PDF text",
                )
        except Exception as e:
            logger.warning("PDF native extraction failed: %s", e)

        if use_ocr and self.textract_client:
            if not page_texts or total_chars < MIN_TEXT_CONTENT_THRESHOLD:
                # No/insufficient native text — full document OCR
                try:
                    ocr_texts = await self._extract_pdf_ocr(await _bytes())
                    if ocr_texts:
                        if page_texts:
                            ocr_dict = {pn: t for t, pn in ocr_texts}
                            merged = []
                            for text, pn in page_texts:
                                ocr = ocr_dict.get(pn, "")
                                if ocr:
                                    best = (
                                        text
                                        if score_extraction_candidate(text)
                                        >= score_extraction_candidate(ocr)
                                        else ocr
                                    )
                                    merged.append((best, pn))
                                else:
                                    merged.append((text, pn))
                            existing = {pn for _, pn in page_texts}
                            for t, pn in ocr_texts:
                                if pn not in existing:
                                    merged.append((t, pn))
                            page_texts = sorted(merged, key=lambda x: x[1])
                        else:
                            page_texts = ocr_texts
                except Exception as e:
                    _log_textract_failure("PDF full OCR", e)
            else:
                # Selective OCR: only pages where native text is weak
                weak_pages = [
                    pn for text, pn in page_texts
                    if is_native_text_probably_weak(text)
                ]
                if weak_pages:
                    logger.info(
                        "_extract_pdf_text: %d/%d weak pages — running selective OCR",
                        len(weak_pages), len(page_texts),
                    )
                    try:
                        ocr_texts = await self._extract_pdf_ocr_for_pages(
                            await _bytes(), weak_pages
                        )
                        if ocr_texts:
                            ocr_dict = {pn: t for t, pn in ocr_texts}
                            merged = []
                            for text, pn in page_texts:
                                ocr = ocr_dict.get(pn, "")
                                if ocr:
                                    best = (
                                        text
                                        if score_extraction_candidate(text)
                                        >= score_extraction_candidate(ocr)
                                        else ocr
                                    )
                                    merged.append((best, pn))
                                else:
                                    merged.append((text, pn))
                            existing = {pn for _, pn in page_texts}
                            for t, pn in ocr_texts:
                                if pn not in existing:
                                    merged.append((t, pn))
                            page_texts = sorted(merged, key=lambda x: x[1])
                    except Exception as e:
                        _log_textract_failure("PDF selective OCR", e)

        if not page_texts:
            return None

        combined = " ".join(t for t, _ in page_texts)
        if not _is_text_meaningful(combined) and self.vision_client:
            vision = (
                await self._caption_pdf_pages_path(file_path)
                if file_path
                else await self._caption_pdf_pages(await _bytes())
            )
            if vision:
                return vision

        # Docling enrichment: append per-page markdown when VECTOR_PDF_DOCLING=true
        if VECTOR_PDF_DOCLING and page_texts is not None:
            dl = await asyncio.to_thread(_extract_docling_per_page_markdown, await _bytes())
            if dl:
                merged = []
                for text, pn in page_texts:
                    extra = (dl.get(pn) or "").strip()
                    if extra:
                        merged.append((f"{text}\n\n[docling]\n{extra}", pn))
                    else:
                        merged.append((text, pn))
                # Add any pages Docling found that native extraction missed
                existing_pages = {pn for _, pn in page_texts}
                for pn, extra in sorted(dl.items()):
                    if pn not in existing_pages:
                        merged.append((f"[docling]\n{extra}", pn))
                page_texts = sorted(merged, key=lambda x: x[1])

        return page_texts or None

    async def _extract_pdf_ocr(self, file_bytes: bytes) -> Optional[List[Tuple[str, int]]]:
        if not self.textract_client:
            return None
        TEXTRACT_LIMIT = 10 * 1024 * 1024
        try:
            if HAS_PYMUPDF:
                doc = await asyncio.to_thread(fitz.open, stream=file_bytes, filetype="pdf")
                total_pages = doc.page_count
                doc.close()
            elif HAS_PYPDF2 and PdfReader is not None:
                reader = await asyncio.to_thread(PdfReader, BytesIO(file_bytes))
                total_pages = len(reader.pages)
            else:
                logger.warning("_extract_pdf_ocr: no PyMuPDF/PyPDF2 — cannot OCR PDF")
                return None

            if len(file_bytes) > TEXTRACT_LIMIT:
                return await self._extract_pdf_ocr_chunked(file_bytes, total_pages)

            response = await asyncio.to_thread(
                self.textract_client.detect_document_text,
                Document={"Bytes": file_bytes},
            )
            page_dict: Dict[int, List[str]] = {}
            for block in response.get("Blocks", []):
                if block["BlockType"] == "LINE":
                    pn = block.get("Page", 1)
                    page_dict.setdefault(pn, []).append(block["Text"])
            result = [("\n".join(lines), pn) for pn, lines in sorted(page_dict.items()) if lines]
            return result if result else await self._extract_pdf_ocr_page_by_page(file_bytes, total_pages)
        except Exception as e:
            _log_textract_failure("_extract_pdf_ocr", e)
            return None

    async def _extract_pdf_ocr_chunked(self, file_bytes: bytes, total_pages: int) -> Optional[List[Tuple[str, int]]]:
        LIMIT = 10 * 1024 * 1024
        bpp = len(file_bytes) / max(total_pages, 1)
        ppc = max(1, int(LIMIT * 0.85 / bpp))
        results = []

        def _build(f, t):
            src = fitz.open(stream=file_bytes, filetype="pdf")
            dst = fitz.open()
            dst.insert_pdf(src, from_page=f, to_page=t)
            buf = BytesIO()
            dst.save(buf, garbage=4, deflate=True)
            data = buf.getvalue()
            dst.close()
            src.close()
            return data

        for start in range(0, total_pages, ppc):
            end = min(start + ppc - 1, total_pages - 1)
            try:
                chunk = await asyncio.to_thread(_build, start, end)
            except Exception as e:
                logger.warning("Chunk build failed %d-%d: %s", start, end, e)
                continue
            if len(chunk) > LIMIT:
                fallback = await self._extract_pdf_ocr_page_by_page(chunk, end - start + 1)
                if fallback:
                    results.extend((t, p + start) for t, p in fallback)
                continue
            try:
                resp = await asyncio.to_thread(
                    self.textract_client.detect_document_text,
                    Document={"Bytes": chunk},
                )
                pd: Dict[int, List[str]] = {}
                for block in resp.get("Blocks", []):
                    if block["BlockType"] == "LINE":
                        pd.setdefault(block.get("Page", 1), []).append(block["Text"])
                for lp, lines in sorted(pd.items()):
                    results.append(("\n".join(lines), start + lp))
            except Exception as e:
                logger.warning("Textract chunk failed: %s", e)
        return results or None

    async def _extract_pdf_ocr_page_by_page(self, file_bytes: bytes, total_pages: int) -> Optional[List[Tuple[str, int]]]:
        if not self.textract_client or not HAS_PYMUPDF:
            return None
        Image, _ImageOps = _pil()
        sem = asyncio.Semaphore(OCR_SEMAPHORE)
        results = []

        async def _process_page(page_num: int):
            async with sem:
                try:
                    def _render():
                        doc = fitz.open(stream=file_bytes, filetype="pdf")
                        pix = doc[page_num].get_pixmap(dpi=200)
                        png = pix.tobytes("png")
                        doc.close()
                        return png

                    png = await asyncio.to_thread(_render)
                    img = await asyncio.to_thread(Image.open, BytesIO(png))
                    img = await asyncio.to_thread(_preprocess_pil_image_for_ocr, img)
                    buf = BytesIO()
                    await asyncio.to_thread(img.save, buf, "JPEG", quality=90)
                    jpeg = buf.getvalue()
                    img.close()
                    buf.close()
                    resp = await asyncio.to_thread(
                        self.textract_client.detect_document_text,
                        Document={"Bytes": jpeg},
                    )
                    text = "\n".join(b["Text"] for b in resp.get("Blocks", []) if b["BlockType"] == "LINE")
                    return (normalize_text(text), page_num + 1) if text.strip() else None
                except Exception as e:
                    logger.warning("Page %d OCR failed: %s", page_num + 1, e)
                    return None

        for r in await asyncio.gather(*[_process_page(i) for i in range(total_pages)]):
            if r:
                results.append(r)
        return sorted(results, key=lambda x: x[1]) or None

    async def _extract_pdf_ocr_for_pages(
        self, file_bytes: bytes, page_numbers: List[int]
    ) -> Optional[List[Tuple[str, int]]]:
        """
        OCR only the given 1-based page numbers (selective weak-page OCR).
        Uses _preprocess_pil_image_for_ocr for better Textract accuracy on
        technical drawings and low-contrast scans.
        """
        if not self.textract_client or not HAS_PYMUPDF:
            return None
        Image, _ImageOps = _pil()
        sem = asyncio.Semaphore(OCR_SEMAPHORE)
        results = []

        async def _process_page(pn: int):
            async with sem:
                try:
                    def _render():
                        doc = fitz.open(stream=file_bytes, filetype="pdf")
                        pix = doc[pn - 1].get_pixmap(dpi=200)
                        png = pix.tobytes("png")
                        doc.close()
                        return png

                    png = await asyncio.to_thread(_render)
                    img = await asyncio.to_thread(Image.open, BytesIO(png))
                    img = await asyncio.to_thread(_preprocess_pil_image_for_ocr, img)
                    buf = BytesIO()
                    await asyncio.to_thread(img.save, buf, "JPEG", quality=90)
                    jpeg = buf.getvalue()
                    img.close()
                    buf.close()
                    resp = await asyncio.to_thread(
                        self.textract_client.detect_document_text,
                        Document={"Bytes": jpeg},
                    )
                    text = "\n".join(
                        b["Text"] for b in resp.get("Blocks", [])
                        if b["BlockType"] == "LINE"
                    )
                    return (normalize_text(text), pn) if text.strip() else None
                except Exception as e:
                    logger.warning("Selective OCR page %d failed: %s", pn, e)
                    return None

        for r in await asyncio.gather(*[_process_page(pn) for pn in page_numbers]):
            if r:
                results.append(r)
        return sorted(results, key=lambda x: x[1]) or None

    # ── Image / TIFF / BIN ─────────────────────────────────────────────────────

    async def _extract_image_text(self, file_bytes=None, file_path=None, file_name: str = ""):
        from pathlib import Path

        if file_bytes is None and file_path:
            file_bytes = await asyncio.to_thread(lambda: open(file_path, "rb").read())
        if not file_bytes:
            return [("[image_file]\nNo readable bytes.", 1)]

        name = (file_name or (Path(file_path).name if file_path else "")).strip()

        ocr_text, caption_text = await asyncio.gather(
            self._textract_image_text(file_bytes),
            self._caption_image(file_bytes, file_name=name),
            return_exceptions=True,
        )
        if isinstance(ocr_text, Exception):
            ocr_text = None
        if isinstance(caption_text, Exception):
            caption_text = None

        parts: List[str] = []
        if name:
            parts.append(f"[image_file] {name}")
        else:
            parts.append("[image_file]")
        if (ocr_text or "").strip():
            parts.append(f"[image_ocr]\n{str(ocr_text).strip()}")
        if (caption_text or "").strip():
            parts.append(f"[image_semantic]\n{str(caption_text).strip()}")

        has_body = (ocr_text or "").strip() or (caption_text or "").strip()
        if has_body:
            return [("\n\n".join(parts), 1)]

        Image, _ImageOps = _pil()
        try:
            img = Image.open(BytesIO(file_bytes))
            fallback = (
                f"[image_file] {name}\n"
                f"Format: {img.format}  Dimensions: {img.size[0]}x{img.size[1]}  "
                f"Size: {len(file_bytes):,} bytes"
            )
            img.close()
        except Exception:
            fallback = f"[image_file] {name}\nSize: {len(file_bytes):,} bytes"
        return [(fallback, 1)]

    async def _textract_image_text(self, file_bytes: bytes) -> Optional[str]:
        if not self.textract_client:
            return None
        Image, _ImageOps = _pil()
        try:
            img = await asyncio.to_thread(Image.open, BytesIO(file_bytes))
            img = await asyncio.to_thread(_preprocess_pil_image_for_ocr, img)
            buf = BytesIO()
            await asyncio.to_thread(img.save, buf, "JPEG", quality=90)
            buf.seek(0)
            resp = await asyncio.to_thread(
                self.textract_client.detect_document_text,
                Document={"Bytes": buf.getvalue()},
            )
            text = "\n".join(b["Text"] for b in resp.get("Blocks", []) if b["BlockType"] == "LINE")
            img.close()
            buf.close()
            return normalize_text(text) if text.strip() else None
        except Exception as e:
            logger.error("Textract image OCR failed: %s", e)
            return None

    async def _caption_image(
        self, file_bytes: bytes, file_name: str = ""
    ) -> Optional[str]:
        """
        Vision LLM caption: outputs DESCRIPTION + TAGS blocks so the combined
        ``content`` string carries both narrative and keyword signal for k-NN
        embedding (mirrors vector.py _enrich_image_with_vision_caption).
        """
        if not self.vision_client:
            return None
        import base64

        Image, _ImageOps = _pil()
        try:
            raw = file_bytes
            # Shrink huge rasters before base64 to stay within API payload limits
            if len(raw) > 6 * 1024 * 1024:
                try:
                    im = await asyncio.to_thread(Image.open, BytesIO(raw))
                    im.thumbnail((2048, 2048), Image.Resampling.LANCZOS)
                    if im.mode in ("RGBA", "P"):
                        im = await asyncio.to_thread(im.convert, "RGB")
                    buf = BytesIO()
                    await asyncio.to_thread(im.save, buf, "JPEG", quality=88)
                    raw = buf.getvalue()
                    im.close()
                    buf.close()
                except Exception as ez:
                    logger.warning("Could not resize large image for vision: %s", ez)

            img = await asyncio.to_thread(Image.open, BytesIO(raw))
            if img.mode != "RGB":
                img = await asyncio.to_thread(img.convert, "RGB")
            buf = BytesIO()
            await asyncio.to_thread(img.save, buf, "JPEG", quality=90)
            b64 = base64.b64encode(buf.getvalue()).decode()
            img.close()
            buf.close()

            name = (file_name or "image").strip() or "image"
            prompt = (
                "You help a document search index. Output exactly two blocks:\n"
                "DESCRIPTION: 2-4 sentences on visible content, setting, subjects, "
                "and any readable text.\n"
                "TAGS: comma-separated list of 12-24 short keywords (objects, actions, "
                "scene, colors) for semantic retrieval.\n"
                f"Filename: {name}\n"
                "If the image is ambiguous, still provide best-effort DESCRIPTION and TAGS."
            )
            for model in ("gpt-4o", "gpt-4o-mini"):
                try:
                    resp = await self.vision_client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
                        ]}],
                        max_tokens=700,
                        temperature=0.2,
                    )
                    content = resp.choices[0].message.content if resp.choices else None
                    if content and content.strip():
                        return content.strip()
                except Exception as e:
                    logger.warning("Caption %s failed: %s", model, e)
            return None
        except Exception as e:
            logger.warning("_caption_image failed: %s", e)
            return None

    async def _caption_pdf_pages(self, file_bytes: bytes, max_pages: int = 10) -> Optional[List[Tuple[str, int]]]:
        if not self.vision_client or not HAS_PYMUPDF:
            return None
        def _render(b, mp):
            doc = fitz.open(stream=b, filetype="pdf")
            pages = []
            for i in range(min(doc.page_count, mp)):
                pix = doc[i].get_pixmap(dpi=250)
                pages.append((pix.tobytes("png"), i + 1))
                pix = None
            doc.close()
            return pages
        rendered = await asyncio.to_thread(_render, file_bytes, max_pages)
        results = []
        for png, pn in rendered:
            cap = await self._caption_image(png)
            if cap:
                results.append((f"[VISION_CAPTION]\n{cap}", pn))
        return results or None

    async def _caption_pdf_pages_path(self, file_path: str, max_pages: int = 10) -> Optional[List[Tuple[str, int]]]:
        if not self.vision_client or not HAS_PYMUPDF:
            return None
        def _render(path, mp):
            doc = fitz.open(path)
            pages = []
            for i in range(min(doc.page_count, mp)):
                pix = doc[i].get_pixmap(dpi=250)
                pages.append((pix.tobytes("png"), i + 1))
                pix = None
            doc.close()
            return pages
        rendered = await asyncio.to_thread(_render, file_path, max_pages)
        results = []
        for png, pn in rendered:
            cap = await self._caption_image(png)
            if cap:
                results.append((f"[VISION_CAPTION]\n{cap}", pn))
        return results or None

    async def _extract_tiff_text(self, file_bytes=None, file_path=None):
        if file_path and file_bytes is None:
            file_bytes = await asyncio.to_thread(lambda: open(file_path, "rb").read())
        if not file_bytes:
            return None
        Image, _ImageOps = _pil()
        img = await asyncio.to_thread(Image.open, BytesIO(file_bytes))
        try:
            if getattr(img, "n_frames", 1) > 1:
                results = []
                for i in range(img.n_frames):
                    await asyncio.to_thread(img.seek, i)
                    frame = await asyncio.to_thread(img.convert, "RGB")
                    buf = BytesIO()
                    await asyncio.to_thread(frame.save, buf, "JPEG", quality=95)
                    r = await self._extract_image_text(buf.getvalue())
                    if r:
                        results.extend(r)
                    frame.close()
                    buf.close()
                return results or None
            else:
                buf = BytesIO()
                rgb = await asyncio.to_thread(img.convert, "RGB")
                await asyncio.to_thread(rgb.save, buf, "JPEG", quality=95)
                return await self._extract_image_text(buf.getvalue())
        finally:
            img.close()
            import gc

            gc.collect()

    async def _extract_bin_text(self, file_bytes=None, file_path=None):
        if file_path and file_bytes is None:
            file_bytes = await asyncio.to_thread(lambda: open(file_path, "rb").read())
        if not file_bytes:
            return None
        Image, _ImageOps = _pil()
        img = await asyncio.to_thread(Image.open, BytesIO(file_bytes))
        if img.mode != "L":
            img = await asyncio.to_thread(img.convert, "L")
        buf = BytesIO()
        await asyncio.to_thread(img.save, buf, "TIFF", compression="tiff_lzw")
        img.close()
        return await self._extract_tiff_text(buf.getvalue())

    # ── Office formats ─────────────────────────────────────────────────────────

    async def _extract_docx_text(self, file_bytes=None, file_path=None):
        from docx import Document as _Docx
        try:
            doc = await asyncio.to_thread(_Docx, file_path if file_path else BytesIO(file_bytes))
            texts: List[str] = []

            # Headers and footers first (section metadata)
            texts.extend(_docx_collect_headers_footers(doc))

            for p in doc.paragraphs:
                if p.text.strip():
                    texts.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
            text = "\n".join(texts)
            if text.strip() and len(text.strip()) >= MIN_TEXT_CONTENT_THRESHOLD:
                return [(text, 1)]
            return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)
        except Exception as e:
            logger.error("DOCX extraction failed: %s", e)
            return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)

    async def _extract_doc_text(self, file_bytes=None, file_path=None):
        if not _is_antiword_available():
            return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)
        import tempfile

        tmp = None
        own_tmp = False
        try:
            if file_path:
                tmp = file_path
            else:
                fd, tmp = tempfile.mkstemp(suffix=".doc")
                os.close(fd)
                with open(tmp, "wb") as f:
                    f.write(file_bytes)
                own_tmp = True
            proc = await asyncio.create_subprocess_exec(
                "antiword", tmp,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
            text = stdout.decode("utf-8", errors="ignore")
            if text.strip() and len(text.strip()) >= MIN_TEXT_CONTENT_THRESHOLD:
                return [(text, 1)]
            return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)
        except Exception as e:
            logger.error("DOC extraction failed: %s", e)
            return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)
        finally:
            if own_tmp and tmp and os.path.exists(tmp):
                os.unlink(tmp)

    async def _extract_excel_text(self, file_bytes=None, file_path=None):
        import openpyxl

        def _sheet_to_text(sheet) -> str:
            title = sheet.title or "Sheet"
            state = getattr(sheet, "sheet_state", None) or "visible"
            header = (
                f"--- Sheet: {title} [{state}] ---"
                if state != "visible"
                else f"--- Sheet: {title} ---"
            )
            row_lines: List[str] = [header]
            for row in sheet.iter_rows(values_only=True):
                cells: List[str] = []
                for cell in row:
                    if cell is None:
                        cells.append("")
                    else:
                        s = str(cell).replace("\t", " ").replace("\r", " ").replace("\n", " ").strip()
                        cells.append(s)
                while cells and cells[-1] == "":
                    cells.pop()
                if not any(cells):
                    continue
                row_lines.append("\t".join(cells))
            return "\n".join(row_lines)

        try:
            wb = await asyncio.to_thread(
                openpyxl.load_workbook,
                file_path if file_path else BytesIO(file_bytes),
                data_only=True,
            )
            preamble = (
                "[Spreadsheet: cell values only; data_only=True. "
                "Formulas appear as last cached values if present.]\n\n"
            )
            blocks = [_sheet_to_text(sheet) for sheet in wb.worksheets]
            wb.close()
            text = preamble + "\n\n".join(blocks)
            if text.strip():
                return [(text, 1)]
            return None
        except Exception as e:
            logger.error("Excel extraction failed: %s", e)
            return None

    async def _extract_csv_text(self, file_bytes=None, file_path=None) -> Optional[List[Tuple[str, int]]]:
        """Parse CSV with automatic encoding + dialect detection (mirrors vector.py)."""
        import csv

        try:
            raw = (
                await asyncio.to_thread(lambda: open(file_path, "rb").read())
                if file_path
                else (file_bytes or b"")
            )
            decoded: Optional[str] = None
            used_enc = "utf-8"
            for enc in ("utf-8-sig", "utf-8", "latin-1"):
                try:
                    decoded = raw.decode(enc)
                    used_enc = enc
                    break
                except UnicodeDecodeError:
                    continue
            if decoded is None:
                decoded = raw.decode("utf-8", errors="replace")
                used_enc = "utf-8 (replace)"

            sample = decoded[: min(len(decoded), 8192)]
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
            except Exception:
                dialect = csv.excel
            f = StringIO(decoded)
            reader = csv.reader(f, dialect)
            rows = list(reader)
            lines: List[str] = [
                f"[CSV encoding={used_enc}, delimiter={repr(dialect.delimiter)}, rows={len(rows)}]"
            ]
            for ri, row in enumerate(rows, start=1):
                clean = [(c or "").replace("\n", " ").strip() for c in row]
                if ri == 1:
                    lines.append("[Row 1 — column names if present] " + "\t".join(clean))
                else:
                    lines.append("\t".join(clean))
            text = "\n".join(lines)
            return [(text, 1)] if text.strip() else None
        except Exception as e:
            logger.error("CSV extraction failed: %s", e)
            return None

    async def _extract_ppt_text(self, file_bytes=None, file_path=None):
        from pptx import Presentation
        try:
            prs = await asyncio.to_thread(Presentation, file_path if file_path else BytesIO(file_bytes))
            texts = []
            for si, slide in enumerate(prs.slides, 1):
                lines: List[str] = [f"--- Slide {si} ---"]
                for shape in slide.shapes:
                    _pptx_collect_shape_text(shape, lines)
                # Speaker notes
                try:
                    if slide.has_notes_slide and slide.notes_slide is not None:
                        nf = slide.notes_slide.notes_text_frame
                        if nf is not None:
                            notes = (nf.text or "").strip()
                            if notes:
                                lines.append("[Speaker notes]")
                                lines.append(notes)
                except Exception:
                    pass
                slide_text = "\n".join(lines)
                if any(ln.strip() for ln in lines[1:]):  # skip header-only slides
                    texts.append((slide_text, si))
            return texts if texts else await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)
        except Exception as e:
            logger.error("PPT extraction failed: %s", e)
            return await self._extract_image_text(file_bytes=file_bytes, file_path=file_path)

    # ── Chunking + Embeddings ──────────────────────────────────────────────────

    async def _chunk_content(self, extracted: List[Tuple[str, int]], file_id: str) -> List[Dict]:
        try:
            from langchain_core.documents import Document as LCDocument

            docs = [LCDocument(page_content=t, metadata={"page_number": pn}) for t, pn in extracted if t.strip()]
            split = self.text_splitter.split_documents(docs)
            return [{"content": d.page_content, "page_number": d.metadata.get("page_number", 1)} for d in split]
        except Exception as e:
            logger.error("_chunk_content error: %s", e)
            return []

    async def _generate_embeddings(self, chunks: List[Dict]) -> List[Optional[List[float]]]:
        texts = [_truncate_chunk(c["content"].strip()) for c in chunks if _chunk_validator(c.get("content", ""))]
        if not texts:
            logger.debug("_generate_embeddings: no valid texts after filtering — returning empty")
            return []
        batches = _create_token_batches(texts)
        logger.debug(
            "_generate_embeddings: %d texts → %d token-batches (max %d tokens each)",
            len(texts), len(batches), MAX_TOKENS_PER_BATCH,
        )
        sem = asyncio.Semaphore(MAX_CONCURRENT_EMBEDDING_BATCHES)

        async def _embed(idx, batch):
            async with sem:
                for attempt in range(3):
                    _t = time.monotonic()
                    try:
                        result = await self.embeddings.aembed_documents(batch)
                        logger.debug(
                            "_generate_embeddings: token-batch %d OK | texts=%d | attempt=%d | elapsed=%.3fs",
                            idx, len(batch), attempt + 1, time.monotonic() - _t,
                        )
                        return result
                    except Exception as e:
                        err = str(e).lower()
                        if attempt < 2 and ("rate limit" in err or "timeout" in err):
                            wait = 60 * (attempt + 1) if "rate limit" in err else 4 ** attempt
                            logger.warning(
                                "_generate_embeddings: token-batch %d attempt %d RETRYABLE error — "
                                "waiting %.0fs | error=%s",
                                idx, attempt + 1, wait, e,
                            )
                            await asyncio.sleep(wait)
                        else:
                            logger.error(
                                "_generate_embeddings: token-batch %d FAILED (attempt %d) | error=%s",
                                idx, attempt + 1, e,
                            )
                            return None
            return None

        results = await asyncio.gather(*[_embed(i, b) for i, b in enumerate(batches)])
        all_emb: List[Optional[List[float]]] = []
        for b, r in zip(batches, results):
            all_emb.extend(r if r else [None] * len(b))
        success = sum(1 for e in all_emb if e is not None)
        logger.debug("_generate_embeddings: %d/%d embeddings succeeded", success, len(all_emb))
        return all_emb


# =============================================================================
# CORE PIPELINE
# =============================================================================

async def _run_vectorize(
    file_doc: Dict[str, Any],
    local_file_path: Optional[str],
    replace: bool,
    use_ocr: bool = True,
) -> None:
    _t_vec = time.monotonic()

    file_id = str(file_doc.get("_id"))
    file_type = (file_doc.get("file_type") or "").lower().strip(".")
    s3_key = file_doc.get("s3_key")

    logger.info(
        "_run_vectorize: START | file_id=%s | type=%s | replace=%s | use_ocr=%s | local_path=%s",
        file_id, file_type, replace, use_ocr, local_file_path,
    )

    vec_svc = StandaloneVectorizationService()
    os_svc = StandaloneOpenSearchService()

    skip_types = [t.strip() for t in (_get_param("SKIP_FILE_TYPES", "") or "").split(",") if t.strip()]
    if file_type in skip_types:
        logger.info(
            "_run_vectorize: file_type='%s' is in SKIP_FILE_TYPES — indexing as unsupported | file_id=%s",
            file_type, file_id,
        )
        await os_svc.index_unsupported_file(file_doc)
        return

    if not replace:
        logger.debug("_run_vectorize: checking if %s is already indexed", file_id)
        already = await os_svc.is_file_indexed(file_id, file_doc.get("full_path", ""))
        if already:
            logger.info("_run_vectorize: %s already indexed — skipping (use replace=True to re-index)", file_id)
            return
        logger.debug("_run_vectorize: %s not yet indexed, proceeding", file_id)

    owns_temp_path = False
    extract_path: Optional[str] = local_file_path

    if extract_path and os.path.isfile(extract_path):
        logger.info(
            "_run_vectorize: using provided local path=%s | size=%.2f MB",
            extract_path, os.path.getsize(extract_path) / (1024 * 1024),
        )
    elif s3_key:
        logger.info("_run_vectorize: no local path — downloading s3_key=%s for file_id=%s", s3_key, file_id)
        try:
            _t = time.monotonic()
            extract_path = download_s3_to_temp(s3_key, file_type)
            owns_temp_path = True
            logger.info(
                "_run_vectorize: S3 download complete | path=%s | size=%.2f MB | elapsed=%.3fs",
                extract_path, os.path.getsize(extract_path) / (1024 * 1024), time.monotonic() - _t,
            )
        except Exception as e:
            logger.error("_run_vectorize: S3 download FAILED | file_id=%s | s3_key=%s | error=%s", file_id, s3_key, e)
            raise
    else:
        logger.warning("_run_vectorize: file_id=%s has no local path and no s3_key — cannot vectorize", file_id)
        return

    try:
        if replace:
            logger.info("_run_vectorize: replace=True — deleting existing OpenSearch docs for %s", file_id)
            _t = time.monotonic()
            await os_svc.delete_file(file_id)
            logger.info("_run_vectorize: existing docs deleted | elapsed=%.3fs", time.monotonic() - _t)

        total_chunks = 0
        batch_num = 0
        first_batch = True

        async for batch in vec_svc.vectorize_file_batched(
            file_id=file_id,
            file_doc=file_doc,
            file_type=file_type,
            file_bytes=None,
            use_ocr=use_ocr,
            file_path=extract_path,
        ):
            chunks = batch.get("chunks", [])
            if not chunks:
                logger.debug("_run_vectorize: empty batch received for %s — skipping", file_id)
                continue

            batch_num += 1
            total_chunks += len(chunks)
            logger.info(
                "_run_vectorize: indexing batch #%d | chunks=%d | cumulative=%d | file_id=%s",
                batch_num, len(chunks), total_chunks, file_id,
            )
            _t = time.monotonic()
            await os_svc.index_file(
                file_id=file_id,
                file_doc=file_doc,
                chunks=chunks,
                replace=(first_batch and replace),
            )
            logger.info(
                "_run_vectorize: batch #%d indexed | elapsed=%.3fs",
                batch_num, time.monotonic() - _t,
            )
            first_batch = False

        if first_batch:
            logger.warning(
                "_run_vectorize: NO CHUNKS produced | file_id=%s | type=%s | path=%s | "
                "File may be empty, image-only, encrypted, or unsupported.",
                file_id, file_type, extract_path,
            )
        else:
            logger.info(
                "_run_vectorize: COMPLETE | file_id=%s | total_chunks=%d | batches=%d | elapsed=%.3fs",
                file_id, total_chunks, batch_num, time.monotonic() - _t_vec,
            )
    finally:
        if owns_temp_path and extract_path:
            _unlink_quiet(extract_path)
            logger.debug("_run_vectorize: cleaned up temp file %s", extract_path)


async def run_ocr_pipeline(
    job: Dict[str, Any],
    db_col,
    verify_task_fn=None,
) -> Dict[str, Any]:
    from bson import ObjectId

    _t_pipeline = time.monotonic()

    file_id = job["file_id"]

    s3_key = (job.get("s3_key") or "").strip()

    file_ext = (job.get("file_ext") or "").lower().lstrip(".")
    collection_name = job.get("collection_name", INSTANT_UPLOAD_COLLECTION)
    replace = job.get("replace", False)
    local_path: Optional[str] = None

    logger.info(
        "run_ocr_pipeline: START | file_id=%s | s3_key=%s | ext=%s | "
        "collection=%s | replace=%s | file_size=%s bytes",
        file_id, s3_key, file_ext, collection_name, replace,
        job.get("file_size", "unknown"),
    )

    # Resolved primary key on the file document (ObjectId or str); used for updates.
    _db_file_pk: Any = None

    def _set_status(status: str, error: Optional[str] = None):
        update: Dict[str, Any] = {"upload_status": status, "updated_at": datetime.now()}
        if error:
            update["index_error"] = error
        pk = _db_file_pk
        if pk is None:
            try:
                pk = ObjectId(file_id)
            except Exception:
                pk = file_id
        try:
            db_col.update_one({"_id": pk}, {"$set": update})
            logger.debug("run_ocr_pipeline: DB status → '%s' for %s", status, file_id)
        except Exception as e:
            logger.warning(
                "run_ocr_pipeline: failed to set status '%s' for %s: %s",
                status, file_id, e,
            )

    try:
        # ── Fetch file document ────────────────────────────────────────────
        _t = time.monotonic()
        file_doc = _find_one_by_file_id(db_col, file_id)
        logger.debug("run_ocr_pipeline: DB fetch took %.3fs | found=%s", time.monotonic() - _t, file_doc is not None)

        if not file_doc:
            logger.warning("run_ocr_pipeline: file_doc not found for file_id=%s", file_id)
            return {"status": "skipped", "file_id": file_id, "reason": "not_found"}

        _db_file_pk = file_doc["_id"]

        logger.info(
            "run_ocr_pipeline: file_doc loaded | name=%s | type=%s | size=%s | status=%s",
            file_doc.get("file_name", "?"),
            file_doc.get("file_type", "?"),
            file_doc.get("file_size", "?"),
            file_doc.get("upload_status", "?"),
        )

        if file_doc.get("upload_status") == "cancelled":
            logger.info("run_ocr_pipeline: %s is cancelled — cleaning up and skipping", file_id)
            _cleanup_cancelled(file_doc, collection_name)
            return {"status": "skipped", "file_id": file_id, "reason": "cancelled"}

        if not s3_key:
            logger.warning("run_ocr_pipeline: job missing s3_key | file_id=%s", file_id)
            _set_status("failed", "missing_s3_key")
            return {"status": "failed", "file_id": file_id, "reason": "missing_s3_key"}

        # ── Download from S3 ───────────────────────────────────────────────
        _set_status("indexing")
        logger.info("run_ocr_pipeline: downloading s3://%s/%s", S3_BUCKET_NAME, s3_key)
        _t = time.monotonic()
        local_path = download_s3_to_temp(s3_key, file_ext)
        _dl_size = os.path.getsize(local_path)
        logger.info(
            "run_ocr_pipeline: S3 download complete | local_path=%s | size=%.2f MB | elapsed=%.3fs",
            local_path, _dl_size / (1024 * 1024), time.monotonic() - _t,
        )

        # ── Validations ────────────────────────────────────────────────────
        logger.debug("run_ocr_pipeline: validating page count for ext=%s", file_ext)
        page_err = validate_page_count(file_ext, local_path)
        if page_err:
            _set_status("failed", page_err)
            logger.warning("run_ocr_pipeline: VALIDATION FAIL (pages) | file_id=%s | reason=%s", file_id, page_err)
            return {"status": "failed", "file_id": file_id, "reason": page_err}

        logger.debug("run_ocr_pipeline: validating character count for ext=%s", file_ext)
        char_err = validate_char_count(file_ext, local_path)
        if char_err:
            _set_status("failed", char_err)
            logger.warning("run_ocr_pipeline: VALIDATION FAIL (chars) | file_id=%s | reason=%s", file_id, char_err)
            return {"status": "failed", "file_id": file_id, "reason": char_err}

        logger.debug("run_ocr_pipeline: validations passed for %s", file_id)

        # ── Cancellation re-check ──────────────────────────────────────────
        fresh = _find_one_by_file_id(db_col, file_id)

        if not fresh or fresh.get("upload_status") == "cancelled":
            logger.info("run_ocr_pipeline: %s cancelled before vectorize — aborting", file_id)
            _cleanup_cancelled(fresh or file_doc, collection_name)
            return {"status": "skipped", "file_id": file_id, "reason": "cancelled_before_vectorize"}

        # ── Compression ────────────────────────────────────────────────────
        _size_before = os.path.getsize(local_path)
        local_path = compress_file_if_needed(local_path, file_ext, s3_key, file_id)
        _size_after = os.path.getsize(local_path)
        if _size_after < _size_before:
            logger.info(
                "run_ocr_pipeline: compression applied | %.2f MB → %.2f MB",
                _size_before / (1024 * 1024), _size_after / (1024 * 1024),
            )
        else:
            logger.debug("run_ocr_pipeline: no compression applied for %s", file_id)

        # ── Vectorize ──────────────────────────────────────────────────────
        logger.info("run_ocr_pipeline: starting vectorization for %s", file_id)
        _t = time.monotonic()
        await _run_vectorize(file_doc, local_path, replace)
        logger.info("run_ocr_pipeline: vectorization done | elapsed=%.3fs", time.monotonic() - _t)

        _set_status("indexed")
        if verify_task_fn:
            verify_task_fn(file_id, collection_name)
            logger.debug("run_ocr_pipeline: verification scheduled for %s", file_id)

        _total = time.monotonic() - _t_pipeline
        logger.info("run_ocr_pipeline: COMPLETE | file_id=%s | total_elapsed=%.3fs", file_id, _total)
        return {"status": "success", "file_id": file_id}

    except Exception as e:
        logger.exception(
            "run_ocr_pipeline: UNHANDLED ERROR | file_id=%s | elapsed=%.3fs | error=%s",
            file_id, time.monotonic() - _t_pipeline, e,
        )
        _set_status("failed", str(e))
        return {"status": "failed", "file_id": file_id, "reason": str(e)}

    finally:
        _unlink_quiet(local_path)


def _cleanup_cancelled(file_doc: Dict[str, Any], collection_name: str) -> None:
    from bson import ObjectId

    try:
        s3_key = file_doc.get("s3_key")
        if s3_key:
            _get_s3_client().delete_object(Bucket=S3_BUCKET_NAME, Key=s3_key)
    except Exception as e:
        logger.warning("_cleanup_cancelled: S3 delete failed: %s", e)
    try:
        db = _get_db()
        col = db[collection_name]
        try:
            col.delete_one({"_id": ObjectId(file_doc["_id"])})
        except Exception:
            col.delete_one({"_id": file_doc["_id"]})
    except Exception as e:
        logger.warning("_cleanup_cancelled: MongoDB delete failed: %s", e)


# =============================================================================
# LAMBDA ENTRY POINT
# =============================================================================

def _schedule_verify(file_id: str, collection_name: str) -> None:
    verify_queue = _get_param("OCR_VERIFY_QUEUE_URL")
    if not verify_queue:
        logger.warning("OCR_VERIFY_QUEUE_URL not set; skipping verification scheduling")
        return
    import boto3

    sqs = boto3.client("sqs", region_name=os.environ.get("AWS_REGION", "us-east-1"))
    sqs.send_message(
        QueueUrl=verify_queue,
        MessageBody=json.dumps({"file_id": file_id, "collection_name": collection_name}),
        DelaySeconds=2,
    )
    logger.info("_schedule_verify: queued verification for %s", file_id)


def handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    """
    Lambda entry point. Triggered by SQS.
    Returns batchItemFailures so SQS retries only failed messages.

    A dedicated log file is written to S3 for every invocation:
        s3://<S3_BUCKET_NAME>/<LOG_S3_PREFIX>/YYYY/MM/DD/<timestamp>_<request-id>.log
    LOG_S3_PREFIX defaults to "logs/lambda".
    """
    _t_handler = time.monotonic()

    # ── Per-invocation log setup ───────────────────────────────────────────────
    _request_id = getattr(context, "aws_request_id", None) or str(uuid.uuid4())
    _now = datetime.utcnow()
    _log_prefix = _get_param("LOG_S3_PREFIX", "logs/lambda")
    _log_key = (
        f"{_log_prefix}/{_now.strftime('%Y/%m/%d')}/"
        f"{_now.strftime('%Y%m%dT%H%M%S')}_{_request_id}.log"
    )
    _s3_log_handler.reset()

    records = event.get("Records", [])
    logger.info(
        "handler: ═══ INVOCATION START ═══ | request_id=%s | records=%d | log_key=%s",
        _request_id, len(records), _log_key,
    )
    logger.info(
        "handler: environment | region=%s | bucket=%s | opensearch_index=%s",
        os.environ.get("AWS_REGION", "?"),
        S3_BUCKET_NAME,
        OPENSEARCH_INDEX_NAME,
    )

    db = _get_db()
    batch_failures = []

    try:
        for record in records:
            message_id = record.get("messageId", "unknown")
            try:
                body = json.loads(record["body"])
            except (json.JSONDecodeError, KeyError) as e:
                logger.error(
                    "handler: malformed SQS message | message_id=%s | error=%s | raw_body=%.200s",
                    message_id, e, record.get("body", ""),
                )
                continue

            file_id = body.get("file_id", "unknown")
            collection_name = body.get("collection_name", INSTANT_UPLOAD_COLLECTION)
            file_size = body.get("file_size", 0)
            file_ext = body.get("file_ext", "?")

            logger.info(
                "handler: ── record START | message_id=%s | file_id=%s | "
                "ext=%s | size=%d bytes (%.2f MB) | collection=%s | replace=%s",
                message_id, file_id, file_ext, file_size, file_size / (1024 * 1024),
                collection_name, body.get("replace", False),
            )

            if file_size >= OCR_LARGE_FILE_THRESHOLD_BYTES:
                logger.info(
                    "handler: file_id=%s size=%.2f MB >= %d bytes threshold — dispatching to ECS",
                    file_id, file_size / (1024 * 1024), OCR_LARGE_FILE_THRESHOLD_BYTES,
                )
                dispatched = _dispatch_to_ecs(record["body"], file_id)
                if not dispatched:
                    batch_failures.append({"itemIdentifier": message_id})
                continue

            col = db[collection_name]
            _t_record = time.monotonic()

            try:
                result = asyncio.run(
                    run_ocr_pipeline(
                        job=body,
                        db_col=col,
                        verify_task_fn=_schedule_verify,
                    )
                )
            except RuntimeError as e:
                if "asyncio.run() cannot be called from a running event loop" in str(e):
                    logger.error(
                        "handler: nested event-loop detected for file_id=%s — "
                        "Lambda may already be running an async context.",
                        file_id,
                    )
                    batch_failures.append({"itemIdentifier": message_id})
                    continue
                raise

            _record_elapsed = time.monotonic() - _t_record
            if result["status"] == "failed":
                batch_failures.append({"itemIdentifier": message_id})
                logger.error(
                    "handler: ── record FAILED | message_id=%s | file_id=%s | "
                    "reason=%s | elapsed=%.3fs",
                    message_id, file_id, result.get("reason"), _record_elapsed,
                )
            else:
                logger.info(
                    "handler: ── record OK | message_id=%s | file_id=%s | "
                    "status=%s | elapsed=%.3fs",
                    message_id, file_id, result["status"], _record_elapsed,
                )

        _total_elapsed = time.monotonic() - _t_handler
        logger.info(
            "handler: ═══ INVOCATION END ═══ | request_id=%s | processed=%d | "
            "failures=%d | elapsed=%.3fs | log_key=%s",
            _request_id, len(records), len(batch_failures), _total_elapsed, _log_key,
        )
        return {"batchItemFailures": batch_failures}

    finally:
        # ── Upload log file to S3 ────────────────────────────────────────────
        _bucket = S3_BUCKET_NAME
        if _bucket:
            _uploaded = _s3_log_handler.flush_to_s3(_bucket, _log_key)
            if not _uploaded:
                print(
                    f"[handler] WARNING: log upload skipped (no lines or upload failed) "
                    f"for key={_log_key}",
                    file=sys.stderr,
                )
        else:
            print(
                "[handler] WARNING: S3_BUCKET_NAME not set — log file NOT uploaded.",
                file=sys.stderr,
            )


# =============================================================================
# LOCAL CLI  (python lambda_function.py response.json)
# =============================================================================

def _event_from_saved_json(data: Dict[str, Any]) -> Dict[str, Any]:
    le = data.get("lambda_event")
    if isinstance(le, dict) and le.get("Records"):
        return le
    if data.get("Records"):
        return data
    # bare job dict — wrap into SQS shape
    return {"Records": [{"messageId": "local-test-001", "body": json.dumps(data)}]}


if __name__ == "__main__":
    from pathlib import Path

    # Load repo-root .env so local runs pick up secrets
    _repo_root = Path(__file__).resolve().parents[1]
    _env_file = _repo_root / ".env"
    if _env_file.is_file():
        load_dotenv(_env_file, override=True)

    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s %(levelname)s %(name)s %(message)s",
    )

    _sqs_payload = os.environ.get("SQS_PAYLOAD")
    if _sqs_payload:
        loaded = json.loads(_sqs_payload)
    elif len(sys.argv) > 1:
        with open(sys.argv[1], encoding="utf-8") as f:
            loaded = json.load(f)
    else:
        loaded = json.load(sys.stdin)

    event = _event_from_saved_json(loaded)
    print(json.dumps(handler(event, None), indent=2, default=_json_default))